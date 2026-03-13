"""
Celery tasks — Railway-compatible (API-based AI, no local GPU needed)
  - Transcription  → OpenAI Whisper API
  - Sentiment      → HuggingFace Inference API (free 30k req/month)
  - Storage        → Cloudinary (free 25GB)
"""

import os
import logging
from celery import Celery
from app.core.config import settings

logger = logging.getLogger(__name__)

celery_app = Celery(
    "murex",
    broker=settings.CELERY_BROKER_URL,
    backend=settings.CELERY_RESULT_BACKEND,
)

celery_app.conf.update(
    task_serializer="json",
    result_serializer="json",
    accept_content=["json"],
    timezone="UTC",
    enable_utc=True,
    task_track_started=True,
    task_acks_late=True,
    worker_prefetch_multiplier=1,
    task_routes={
        "app.tasks.tasks.send_survey_email":   {"queue": "notifications"},
        "app.tasks.tasks.send_survey_sms":     {"queue": "notifications"},
        "app.tasks.tasks.transcribe_media":    {"queue": "ai"},
        "app.tasks.tasks.analyze_sentiment":   {"queue": "ai"},
        "app.tasks.tasks.score_ms_submission": {"queue": "scoring"},
        "app.tasks.tasks.generate_nps_report": {"queue": "reports"},
        "app.tasks.tasks.export_to_excel":     {"queue": "reports"},
        "app.tasks.tasks.export_to_pptx":      {"queue": "reports"},
    },
    beat_schedule={
        "close-expired-surveys":      {"task": "app.tasks.tasks.close_expired_surveys",      "schedule": 3600.0},
        "send-assignment-reminders":  {"task": "app.tasks.tasks.send_assignment_reminders",   "schedule": 86400.0},
    },
)


# ─────────────────────────────────────────────────────────
# MODULE A — NOTIFICATIONS
# ─────────────────────────────────────────────────────────

@celery_app.task(bind=True, max_retries=3, default_retry_delay=60)
def send_survey_email(self, distribution_id: int):
    """Send survey invite via SendGrid."""
    try:
        import asyncio
        from sendgrid import SendGridAPIClient
        from sendgrid.helpers.mail import Mail

        async def _send():
            from app.db.session import AsyncSessionLocal
            from sqlalchemy import select
            from app.models.models import Distribution, Survey
            from datetime import datetime, timezone

            async with AsyncSessionLocal() as db:
                dist = (await db.execute(select(Distribution).where(Distribution.id == distribution_id))).scalar_one_or_none()
                if not dist or dist.is_sent:
                    return
                survey = (await db.execute(select(Survey).where(Survey.id == dist.survey_id))).scalar_one_or_none()
                survey_url = f"{settings.FRONTEND_URL}/survey/{dist.token}"
                mail = Mail(
                    from_email=settings.SENDGRID_FROM_EMAIL,
                    to_emails=dist.recipient_email,
                    subject=f"Survey: {survey.title}",
                    html_content=f"""
                    <div style="font-family:Arial,sans-serif;max-width:600px;margin:0 auto">
                      <h2>{survey.title}</h2>
                      <p>You are invited to share your feedback.</p>
                      <a href="{survey_url}"
                         style="background:#2563eb;color:#fff;padding:14px 28px;
                                border-radius:8px;text-decoration:none;display:inline-block">
                        Start Survey →
                      </a>
                    </div>""",
                )
                SendGridAPIClient(settings.SENDGRID_API_KEY).send(mail)
                dist.is_sent = True
                dist.sent_at = datetime.now(timezone.utc)
                await db.commit()

        asyncio.get_event_loop().run_until_complete(_send())
    except Exception as exc:
        logger.error(f"Email failed dist {distribution_id}: {exc}")
        raise self.retry(exc=exc)


@celery_app.task(bind=True, max_retries=3, default_retry_delay=60)
def send_survey_sms(self, distribution_id: int):
    """Send survey invite via Twilio SMS."""
    try:
        import asyncio
        from twilio.rest import Client

        async def _send():
            from app.db.session import AsyncSessionLocal
            from sqlalchemy import select
            from app.models.models import Distribution
            from datetime import datetime, timezone

            async with AsyncSessionLocal() as db:
                dist = (await db.execute(select(Distribution).where(Distribution.id == distribution_id))).scalar_one_or_none()
                if not dist or dist.is_sent or not dist.recipient_phone:
                    return
                Client(settings.TWILIO_ACCOUNT_SID, settings.TWILIO_AUTH_TOKEN).messages.create(
                    body=f"Survey invite: {settings.FRONTEND_URL}/s/{dist.token}",
                    from_=settings.TWILIO_FROM_NUMBER,
                    to=dist.recipient_phone,
                )
                dist.is_sent = True
                dist.sent_at = datetime.now(timezone.utc)
                await db.commit()

        asyncio.get_event_loop().run_until_complete(_send())
    except Exception as exc:
        raise self.retry(exc=exc)


@celery_app.task
def close_expired_surveys():
    import asyncio
    from datetime import datetime, timezone

    async def _run():
        from app.db.session import AsyncSessionLocal
        from sqlalchemy import update
        from app.models.models import Survey, SurveyStatus
        async with AsyncSessionLocal() as db:
            await db.execute(
                update(Survey)
                .where(Survey.end_date < datetime.now(timezone.utc), Survey.status == SurveyStatus.ACTIVE)
                .values(status=SurveyStatus.CLOSED)
            )
            await db.commit()

    asyncio.get_event_loop().run_until_complete(_run())


# ─────────────────────────────────────────────────────────
# MODULE C — AI (OpenAI Whisper API + HuggingFace free API)
# ─────────────────────────────────────────────────────────

@celery_app.task(bind=True, max_retries=2, time_limit=600)
def transcribe_media(self, media_file_id: int):
    """
    Transcribe using OpenAI Whisper API (no local model/GPU needed).
    Cost: ~$0.006/min. A 10-min interview = ~$0.06.
    """
    import asyncio
    import tempfile
    import httpx

    async def _run():
        from app.db.session import AsyncSessionLocal
        from sqlalchemy import select
        from app.models.models import MediaFile, Transcript, MediaStatus
        from openai import OpenAI

        async with AsyncSessionLocal() as db:
            mf = (await db.execute(select(MediaFile).where(MediaFile.id == media_file_id))).scalar_one_or_none()
            if not mf:
                return
            mf.status = MediaStatus.PROCESSING
            await db.commit()

            try:
                # Download from Cloudinary URL
                async with httpx.AsyncClient(timeout=120) as client:
                    resp = await client.get(mf.s3_key)
                    audio_bytes = resp.content

                oai = OpenAI(api_key=settings.OPENAI_API_KEY)
                ext = mf.file_name.rsplit(".", 1)[-1] if "." in mf.file_name else "mp4"

                with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
                    tmp.write(audio_bytes)
                    tmp_path = tmp.name

                with open(tmp_path, "rb") as f:
                    result = oai.audio.transcriptions.create(
                        model="whisper-1",
                        file=f,
                        language=mf.language if mf.language not in ("auto", "") else None,
                        response_format="verbose_json",
                        timestamp_granularities=["segment"],
                    )
                os.unlink(tmp_path)

                segments = [
                    {"start": s.start, "end": s.end, "text": s.text.strip(), "speaker": None}
                    for s in (result.segments or [])
                ]
                db.add(Transcript(
                    media_file_id=mf.id,
                    full_text=result.text,
                    segments=segments,
                    language_detected=result.language,
                ))
                mf.status = MediaStatus.TRANSCRIBED
                await db.commit()
                analyze_sentiment.delay(media_file_id)

            except Exception as e:
                mf.status = MediaStatus.FAILED
                await db.commit()
                raise e

    asyncio.get_event_loop().run_until_complete(_run())


@celery_app.task(bind=True, max_retries=2)
def analyze_sentiment(self, media_file_id: int):
    """
    Arabic/English sentiment via HuggingFace Inference API.
    Free: 30,000 requests/month, no credit card.
    """
    import asyncio
    import httpx
    from collections import Counter
    import re

    async def _run():
        from app.db.session import AsyncSessionLocal
        from sqlalchemy import select
        from app.models.models import MediaFile, Transcript, QualAnalysis, MediaStatus

        async with AsyncSessionLocal() as db:
            mf = (await db.execute(select(MediaFile).where(MediaFile.id == media_file_id))).scalar_one_or_none()
            tr = (await db.execute(select(Transcript).where(Transcript.media_file_id == media_file_id))).scalar_one_or_none()
            if not mf or not tr or not tr.full_text:
                return

            model = ("CAMeL-Lab/bert-base-arabic-camelbert-mix-sentiment"
                     if mf.language in ("ar", "arabic")
                     else "cardiffnlp/twitter-xlm-roberta-base-sentiment")

            hf_url = f"https://api-inference.huggingface.co/models/{model}"
            headers = {"Authorization": f"Bearer {settings.HUGGINGFACE_TOKEN}"}
            words = tr.full_text.split()
            chunks = [" ".join(words[i:i+200]) for i in range(0, min(len(words), 2000), 200)]

            sentiments = []
            async with httpx.AsyncClient(timeout=30) as client:
                for chunk in chunks:
                    if not chunk.strip():
                        continue
                    r = await client.post(hf_url, headers=headers, json={"inputs": chunk})
                    if r.status_code == 200:
                        data = r.json()
                        if data and isinstance(data[0], list):
                            sentiments.append(max(data[0], key=lambda x: x["score"]))

            if sentiments:
                pos = sum(1 for s in sentiments if "pos" in s["label"].lower())
                neg = sum(1 for s in sentiments if "neg" in s["label"].lower())
                overall = "positive" if pos > neg else "negative" if neg > pos else "neutral"
                avg_score = sum(s["score"] for s in sentiments) / len(sentiments)
            else:
                overall, avg_score = "neutral", 0.5

            stop_words = {"the","a","an","and","or","in","on","at","to","of","is","are","was",
                          "من","في","على","إلى","مع","هذا","هذه","أن","كان","لا"}
            tokens = re.findall(r'\b[\w\u0600-\u06FF]{3,}\b', tr.full_text.lower())
            keywords = [{"word": k, "count": v}
                        for k, v in Counter(w for w in tokens if w not in stop_words).most_common(20)]

            db.add(QualAnalysis(
                media_file_id=mf.id,
                overall_sentiment=overall,
                sentiment_score=round(avg_score, 4),
                keywords=keywords,
                topics=[],
                sentiment_timeline=sentiments[:10],
                summary=f"Analyzed {len(sentiments)} segments. Overall: {overall}.",
            ))
            mf.status = MediaStatus.ANALYZED
            await db.commit()

    asyncio.get_event_loop().run_until_complete(_run())


# ─────────────────────────────────────────────────────────
# MODULE B — MYSTERY SHOPPING
# ─────────────────────────────────────────────────────────

@celery_app.task(bind=True, max_retries=2)
def score_ms_submission(self, submission_id: int):
    import asyncio

    async def _run():
        from app.db.session import AsyncSessionLocal
        from sqlalchemy import select
        from app.models.models import MSSubmission, AuditForm

        async with AsyncSessionLocal() as db:
            sub = (await db.execute(select(MSSubmission).where(MSSubmission.id == submission_id))).scalar_one_or_none()
            if not sub:
                return
            form = (await db.execute(select(AuditForm).where(AuditForm.id == sub.form_id))).scalar_one_or_none()
            if not form or not form.sections:
                return

            total, max_score = 0.0, 0.0
            for section in form.sections:
                for q in section.get("questions", []):
                    weight = float(q.get("weight", 1.0))
                    max_score += weight
                    answer = sub.answers.get(str(q["id"]))
                    if answer is not None:
                        if q["type"] == "boolean":
                            total += weight if answer is True else 0
                        elif q["type"] == "rating":
                            total += (float(answer) / q.get("max", 5)) * weight
            sub.score = round((total / max_score) * 100, 2) if max_score else 0
            sub.qa_status = "auto_scored"
            await db.commit()

    asyncio.get_event_loop().run_until_complete(_run())


@celery_app.task
def send_assignment_reminders():
    import asyncio
    from datetime import datetime, timezone, timedelta

    async def _run():
        from app.db.session import AsyncSessionLocal
        from sqlalchemy import select
        from app.models.models import MSAssignment, TaskStatus, ShopperProfile, User
        from twilio.rest import Client

        async with AsyncSessionLocal() as db:
            tomorrow = datetime.now(timezone.utc) + timedelta(days=1)
            assignments = (await db.execute(
                select(MSAssignment).where(
                    MSAssignment.due_date <= tomorrow,
                    MSAssignment.status == TaskStatus.ASSIGNED,
                )
            )).scalars().all()

            tw = Client(settings.TWILIO_ACCOUNT_SID, settings.TWILIO_AUTH_TOKEN)
            for a in assignments:
                if not a.shopper_id:
                    continue
                profile = (await db.execute(select(ShopperProfile).where(ShopperProfile.id == a.shopper_id))).scalar_one_or_none()
                user = (await db.execute(select(User).where(User.id == profile.user_id))).scalar_one_or_none() if profile else None
                if user and user.phone:
                    tw.messages.create(
                        body="Reminder: Mystery shopping visit due tomorrow. Check the app.",
                        from_=settings.TWILIO_FROM_NUMBER,
                        to=user.phone,
                    )

    asyncio.get_event_loop().run_until_complete(_run())


# ─────────────────────────────────────────────────────────
# REPORTING
# ─────────────────────────────────────────────────────────

@celery_app.task
def generate_nps_report(survey_id: int) -> dict:
    import asyncio

    async def _run():
        from app.db.session import AsyncSessionLocal
        from sqlalchemy import select
        from app.models.models import SurveyResponse

        async with AsyncSessionLocal() as db:
            responses = (await db.execute(
                select(SurveyResponse).where(SurveyResponse.survey_id == survey_id, SurveyResponse.is_complete == True)
            )).scalars().all()
            if not responses:
                return {"survey_id": survey_id, "total_responses": 0, "nps_score": None}
            nps = [r.nps_score for r in responses if r.nps_score is not None]
            promoters  = sum(1 for s in nps if s >= 9)
            detractors = sum(1 for s in nps if s <= 6)
            passives   = len(nps) - promoters - detractors
            csat = [r.csat_score for r in responses if r.csat_score is not None]
            return {
                "survey_id": survey_id,
                "total_responses": len(responses),
                "promoters": promoters, "passives": passives, "detractors": detractors,
                "nps_score": round(((promoters - detractors) / len(nps)) * 100, 1) if nps else 0,
                "avg_csat": round(sum(csat)/len(csat), 2) if csat else None,
            }

    return asyncio.get_event_loop().run_until_complete(_run())


@celery_app.task
def export_to_excel(survey_id: int, output_path: str):
    import asyncio, pandas as pd

    async def _run():
        from app.db.session import AsyncSessionLocal
        from sqlalchemy import select
        from app.models.models import SurveyResponse, Answer

        async with AsyncSessionLocal() as db:
            responses = (await db.execute(select(SurveyResponse).where(SurveyResponse.survey_id == survey_id))).scalars().all()
            rows = []
            for r in responses:
                answers = (await db.execute(select(Answer).where(Answer.response_id == r.id))).scalars().all()
                row = {"response_id": r.id, "submitted_at": r.submitted_at, "nps_score": r.nps_score, "csat_score": r.csat_score}
                for a in answers:
                    row[f"q_{a.question_id}"] = a.value_text or a.value_numeric or str(a.value_choices or "")
                rows.append(row)
            pd.DataFrame(rows).to_excel(output_path, index=False, engine="openpyxl")
            return output_path

    return asyncio.get_event_loop().run_until_complete(_run())


@celery_app.task
def export_to_pptx(survey_id: int, output_path: str):
    from pptx import Presentation
    from pptx.util import Inches, Pt

    data = generate_nps_report(survey_id)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Survey NPS Report"
    slide.placeholders[1].text = f"Survey ID: {survey_id}"
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    tf = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5)).text_frame
    p = tf.add_paragraph()
    p.text = f"NPS Score: {data.get('nps_score', 'N/A')}"
    p.font.size = Pt(36)
    p.font.bold = True
    for key in ("total_responses","promoters","passives","detractors","avg_csat"):
        p2 = tf.add_paragraph()
        p2.text = f"{key.replace('_',' ').title()}: {data.get(key,'N/A')}"
        p2.font.size = Pt(18)
    prs.save(output_path)
    return output_path
