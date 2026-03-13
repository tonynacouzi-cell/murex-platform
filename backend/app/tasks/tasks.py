"""
Celery tasks — Railway-compatible version
- Media: Cloudinary (free 25GB) instead of S3/MinIO
- Transcription: OpenAI Whisper API instead of local model
- Sentiment: OpenAI GPT-4o-mini (cheaper than running HuggingFace locally)
"""

import logging
from celery import Celery
from app.core.config import settings

logger = logging.getLogger(__name__)

celery_app = Celery("murex")
celery_app.config_from_object({
    "broker_url": settings.celery_broker,
    "result_backend": settings.celery_backend,
    "task_serializer": "json",
    "result_serializer": "json",
    "accept_content": ["json"],
    "timezone": "UTC",
    "enable_utc": True,
    "task_track_started": True,
    "task_acks_late": True,
    "worker_prefetch_multiplier": 1,
    "task_routes": {
        "app.tasks.tasks.send_survey_email":   {"queue": "notifications"},
        "app.tasks.tasks.send_survey_sms":     {"queue": "notifications"},
        "app.tasks.tasks.send_whatsapp":       {"queue": "notifications"},
        "app.tasks.tasks.transcribe_media":    {"queue": "ai"},
        "app.tasks.tasks.analyze_sentiment":   {"queue": "ai"},
        "app.tasks.tasks.score_ms_submission": {"queue": "scoring"},
        "app.tasks.tasks.generate_nps_report": {"queue": "reports"},
        "app.tasks.tasks.export_to_excel":     {"queue": "reports"},
        "app.tasks.tasks.export_to_pptx":      {"queue": "reports"},
    },
    "beat_schedule": {
        "close-expired-surveys": {
            "task": "app.tasks.tasks.close_expired_surveys",
            "schedule": 3600.0,
        },
        "send-assignment-reminders": {
            "task": "app.tasks.tasks.send_assignment_reminders",
            "schedule": 86400.0,
        },
    },
})


def run_async(coro):
    """Helper: run async functions inside sync Celery tasks."""
    import asyncio
    try:
        loop = asyncio.get_event_loop()
        if loop.is_closed():
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
        return loop.run_until_complete(coro)
    except RuntimeError:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        return loop.run_until_complete(coro)


# ─────────────────────────────────────────────
# MODULE A — NOTIFICATIONS
# ─────────────────────────────────────────────

@celery_app.task(bind=True, max_retries=3, default_retry_delay=60)
def send_survey_email(self, distribution_id: int):
    """Send survey invite via SendGrid."""
    async def _send():
        from sendgrid import SendGridAPIClient
        from sendgrid.helpers.mail import Mail
        from app.db.session import AsyncSessionLocal
        from sqlalchemy import select
        from app.models.models import Distribution, Survey
        from datetime import datetime, timezone

        async with AsyncSessionLocal() as db:
            dist = (await db.execute(select(Distribution).where(Distribution.id == distribution_id))).scalar_one_or_none()
            if not dist or dist.is_sent:
                return
            survey = (await db.execute(select(Survey).where(Survey.id == dist.survey_id))).scalar_one_or_none()
            survey_url = f"{settings.FRONTEND_URL}/survey/{dist.token}"

            html = f"""
            <div style="font-family:Arial,sans-serif;max-width:600px;margin:auto">
              <h2 style="color:#1E3A5F">{survey.title}</h2>
              <p>You're invited to share your feedback. It takes less than 5 minutes.</p>
              <a href="{survey_url}"
                 style="display:inline-block;background:#1E3A5F;color:#fff;padding:14px 28px;
                        border-radius:8px;text-decoration:none;font-size:16px">
                Start Survey →
              </a>
              <p style="color:#999;font-size:12px;margin-top:24px">
                This link is unique to you. Do not share it.
              </p>
            </div>"""

            sg = SendGridAPIClient(settings.SENDGRID_API_KEY)
            sg.send(Mail(
                from_email=settings.SENDGRID_FROM_EMAIL,
                to_emails=dist.recipient_email,
                subject=f"Your feedback matters — {survey.title}",
                html_content=html,
            ))
            dist.is_sent = True
            dist.sent_at = datetime.now(timezone.utc)
            await db.commit()
            logger.info(f"Email sent: distribution {distribution_id}")

    try:
        run_async(_send())
    except Exception as exc:
        logger.error(f"Email failed for dist {distribution_id}: {exc}")
        raise self.retry(exc=exc)


@celery_app.task(bind=True, max_retries=3, default_retry_delay=60)
def send_survey_sms(self, distribution_id: int):
    """Send survey invite via Twilio SMS."""
    async def _send():
        from twilio.rest import Client
        from app.db.session import AsyncSessionLocal
        from sqlalchemy import select
        from app.models.models import Distribution
        from datetime import datetime, timezone

        async with AsyncSessionLocal() as db:
            dist = (await db.execute(select(Distribution).where(Distribution.id == distribution_id))).scalar_one_or_none()
            if not dist or dist.is_sent or not dist.recipient_phone:
                return
            survey_url = f"{settings.FRONTEND_URL}/s/{dist.token}"
            Client(settings.TWILIO_ACCOUNT_SID, settings.TWILIO_AUTH_TOKEN).messages.create(
                body=f"Share your feedback — complete our short survey: {survey_url}",
                from_=settings.TWILIO_FROM_NUMBER,
                to=dist.recipient_phone,
            )
            dist.is_sent = True
            dist.sent_at = datetime.now(timezone.utc)
            await db.commit()

    try:
        run_async(_send())
    except Exception as exc:
        raise self.retry(exc=exc)


@celery_app.task
def send_whatsapp(phone: str, message: str):
    """Send WhatsApp message via Twilio."""
    try:
        from twilio.rest import Client
        Client(settings.TWILIO_ACCOUNT_SID, settings.TWILIO_AUTH_TOKEN).messages.create(
            from_=f"whatsapp:{settings.TWILIO_FROM_NUMBER}",
            to=f"whatsapp:{phone}",
            body=message,
        )
    except Exception as e:
        logger.error(f"WhatsApp failed to {phone}: {e}")
        raise


@celery_app.task
def close_expired_surveys():
    """Periodic: auto-close surveys past end_date."""
    async def _close():
        from app.db.session import AsyncSessionLocal
        from sqlalchemy import update
        from app.models.models import Survey, SurveyStatus
        from datetime import datetime, timezone

        async with AsyncSessionLocal() as db:
            await db.execute(
                update(Survey)
                .where(Survey.end_date < datetime.now(timezone.utc), Survey.status == SurveyStatus.ACTIVE)
                .values(status=SurveyStatus.CLOSED)
            )
            await db.commit()

    run_async(_close())


# ─────────────────────────────────────────────
# MODULE B — SCORING
# ─────────────────────────────────────────────

@celery_app.task(bind=True, max_retries=2)
def score_ms_submission(self, submission_id: int):
    """Auto-score mystery shopping submission."""
    async def _score():
        from app.db.session import AsyncSessionLocal
        from sqlalchemy import select
        from app.models.models import MSSubmission, AuditForm

        async with AsyncSessionLocal() as db:
            sub = (await db.execute(select(MSSubmission).where(MSSubmission.id == submission_id))).scalar_one_or_none()
            if not sub:
                return
            form = (await db.execute(select(AuditForm).where(AuditForm.id == sub.form_id))).scalar_one_or_none()
            if not form or not form.sections:
                return

            total, max_score = 0.0, 0.0
            for section in form.sections:
                for q in section.get("questions", []):
                    weight = float(q.get("weight", 1.0))
                    max_score += weight
                    answer = sub.answers.get(str(q["id"]))
                    if answer is not None:
                        if q["type"] == "boolean":
                            total += weight if answer is True else 0
                        elif q["type"] == "rating":
                            total += (float(answer) / q.get("max", 5)) * weight

            sub.score = round((total / max_score) * 100, 2) if max_score > 0 else 0
            sub.qa_status = "auto_scored"
            await db.commit()
            logger.info(f"Submission {submission_id} scored: {sub.score}")

    try:
        run_async(_score())
    except Exception as exc:
        raise self.retry(exc=exc)


@celery_app.task
def send_assignment_reminders():
    """Daily: remind shoppers of due assignments."""
    async def _remind():
        from app.db.session import AsyncSessionLocal
        from sqlalchemy import select
        from app.models.models import MSAssignment, TaskStatus, ShopperProfile, User
        from datetime import datetime, timezone, timedelta

        async with AsyncSessionLocal() as db:
            tomorrow = datetime.now(timezone.utc) + timedelta(days=1)
            result = await db.execute(
                select(MSAssignment).where(
                    MSAssignment.due_date <= tomorrow,
                    MSAssignment.status == TaskStatus.ASSIGNED,
                )
            )
            for a in result.scalars().all():
                if a.shopper_id:
                    logger.info(f"Reminder queued for assignment {a.id}")

    run_async(_remind())


# ─────────────────────────────────────────────
# MODULE C — AI PIPELINE (API-based, Railway-safe)
# ─────────────────────────────────────────────

@celery_app.task(bind=True, max_retries=2, time_limit=600)
def transcribe_media(self, media_file_id: int):
    """
    Download from Cloudinary, transcribe via OpenAI Whisper API.
    Railway-safe: no local GPU/model required, just API calls.
    """
    async def _transcribe():
        import tempfile, os, httpx
        from openai import OpenAI
        import cloudinary.api
        from app.db.session import AsyncSessionLocal
        from sqlalchemy import select
        from app.models.models import MediaFile, Transcript, MediaStatus

        async with AsyncSessionLocal() as db:
            mf = (await db.execute(select(MediaFile).where(MediaFile.id == media_file_id))).scalar_one_or_none()
            if not mf:
                return

            mf.status = MediaStatus.PROCESSING
            await db.commit()

            try:
                # Download from Cloudinary
                download_url = cloudinary.utils.cloudinary_url(mf.cloudinary_public_id)[0]
                async with httpx.AsyncClient(timeout=120) as client:
                    response = await client.get(download_url)
                    response.raise_for_status()

                with tempfile.NamedTemporaryFile(suffix=".mp4", delete=False) as tmp:
                    tmp.write(response.content)
                    tmp_path = tmp.name

                # OpenAI Whisper API — supports Arabic natively
                client = OpenAI(api_key=settings.OPENAI_API_KEY)
                with open(tmp_path, "rb") as audio_file:
                    result = client.audio.transcriptions.create(
                        model="whisper-1",
                        file=audio_file,
                        language=mf.language if mf.language != "auto" else None,
                        response_format="verbose_json",
                        timestamp_granularities=["segment"],
                    )

                segments = [
                    {"start": s.start, "end": s.end, "text": s.text.strip(), "speaker": None}
                    for s in (result.segments or [])
                ]

                transcript = Transcript(
                    media_file_id=mf.id,
                    full_text=result.text,
                    segments=segments,
                    language_detected=result.language,
                )
                db.add(transcript)
                mf.status = MediaStatus.TRANSCRIBED
                await db.commit()

                analyze_sentiment.delay(media_file_id)
                os.unlink(tmp_path)
                logger.info(f"Transcription done: media {media_file_id}")

            except Exception as e:
                mf.status = MediaStatus.FAILED
                await db.commit()
                raise e

    try:
        run_async(_transcribe())
    except Exception as exc:
        raise self.retry(exc=exc)


@celery_app.task(bind=True, max_retries=2)
def analyze_sentiment(self, media_file_id: int):
    """
    Arabic/English sentiment analysis via GPT-4o-mini.
    Cheaper and Railway-safe vs. running HuggingFace locally.
    """
    async def _analyze():
        import json, re
        from collections import Counter
        from openai import OpenAI
        from app.db.session import AsyncSessionLocal
        from sqlalchemy import select
        from app.models.models import MediaFile, Transcript, QualAnalysis, MediaStatus

        async with AsyncSessionLocal() as db:
            mf = (await db.execute(select(MediaFile).where(MediaFile.id == media_file_id))).scalar_one_or_none()
            tr = (await db.execute(select(Transcript).where(Transcript.media_file_id == media_file_id))).scalar_one_or_none()
            if not mf or not tr or not tr.full_text:
                return

            try:
                client = OpenAI(api_key=settings.OPENAI_API_KEY)

                # Truncate to ~3000 words for cost control
                text_sample = " ".join(tr.full_text.split()[:3000])

                lang_hint = "Arabic" if mf.language in ("ar", "arabic") else "English"
                prompt = f"""Analyze this {lang_hint} transcript and return ONLY a JSON object with:
{{
  "overall_sentiment": "positive" | "negative" | "neutral",
  "sentiment_score": 0.0-1.0,
  "summary": "2-3 sentence summary",
  "keywords": [{{"word": "...", "count": N}}, ...] (top 10),
  "topics": [{{"topic": "...", "relevance": 0.0-1.0}}, ...] (top 5)
}}

Transcript:
{text_sample}"""

                response = client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[{"role": "user", "content": prompt}],
                    response_format={"type": "json_object"},
                    temperature=0.1,
                )

                data = json.loads(response.choices[0].message.content)

                analysis = QualAnalysis(
                    media_file_id=mf.id,
                    overall_sentiment=data.get("overall_sentiment", "neutral"),
                    sentiment_score=data.get("sentiment_score", 0.5),
                    keywords=data.get("keywords", []),
                    topics=data.get("topics", []),
                    summary=data.get("summary", ""),
                    sentiment_timeline=[],
                )
                db.add(analysis)
                mf.status = MediaStatus.ANALYZED
                await db.commit()
                logger.info(f"Analysis done: media {media_file_id}")

            except Exception as e:
                logger.error(f"Analysis failed: media {media_file_id}: {e}")
                raise e

    try:
        run_async(_analyze())
    except Exception as exc:
        raise self.retry(exc=exc)


# ─────────────────────────────────────────────
# REPORTING
# ─────────────────────────────────────────────

@celery_app.task
def generate_nps_report(survey_id: int) -> dict:
    async def _calc():
        from app.db.session import AsyncSessionLocal
        from sqlalchemy import select
        from app.models.models import SurveyResponse

        async with AsyncSessionLocal() as db:
            result = await db.execute(
                select(SurveyResponse).where(
                    SurveyResponse.survey_id == survey_id,
                    SurveyResponse.is_complete == True,
                )
            )
            responses = result.scalars().all()
            if not responses:
                return {"survey_id": survey_id, "total_responses": 0, "nps_score": None}

            nps_scores = [r.nps_score for r in responses if r.nps_score is not None]
            promoters  = sum(1 for s in nps_scores if s >= 9)
            passives   = sum(1 for s in nps_scores if 7 <= s <= 8)
            detractors = sum(1 for s in nps_scores if s <= 6)
            total      = len(nps_scores)
            nps        = round(((promoters - detractors) / total) * 100, 1) if total else 0
            csat_vals  = [r.csat_score for r in responses if r.csat_score is not None]

            return {
                "survey_id": survey_id,
                "total_responses": len(responses),
                "promoters": promoters,
                "passives": passives,
                "detractors": detractors,
                "nps_score": nps,
                "avg_csat": round(sum(csat_vals) / len(csat_vals), 2) if csat_vals else None,
            }

    return run_async(_calc())


@celery_app.task
def export_to_excel(survey_id: int, output_path: str):
    async def _export():
        import pandas as pd
        from app.db.session import AsyncSessionLocal
        from sqlalchemy import select
        from app.models.models import SurveyResponse, Answer

        async with AsyncSessionLocal() as db:
            responses = (await db.execute(
                select(SurveyResponse).where(SurveyResponse.survey_id == survey_id)
            )).scalars().all()

            rows = []
            for r in responses:
                answers = (await db.execute(select(Answer).where(Answer.response_id == r.id))).scalars().all()
                row = {
                    "response_id": r.id,
                    "submitted_at": r.submitted_at,
                    "is_complete": r.is_complete,
                    "nps_score": r.nps_score,
                    "csat_score": r.csat_score,
                }
                for a in answers:
                    row[f"q_{a.question_id}"] = a.value_text or a.value_numeric or str(a.value_choices or "")
                rows.append(row)

            pd.DataFrame(rows).to_excel(output_path, index=False, engine="openpyxl")
        return output_path

    return run_async(_export())


@celery_app.task
def export_to_pptx(survey_id: int, output_path: str):
    from pptx import Presentation
    from pptx.util import Inches, Pt

    nps_data = generate_nps_report(survey_id)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Survey Report"
    slide.placeholders[1].text = f"Survey ID: {survey_id}"

    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    tf = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5)).text_frame
    tf.word_wrap = True

    p = tf.add_paragraph()
    p.text = f"NPS Score: {nps_data.get('nps_score', 'N/A')}"
    p.font.size = Pt(32)
    p.font.bold = True

    for key in ("total_responses", "promoters", "passives", "detractors", "avg_csat"):
        p2 = tf.add_paragraph()
        p2.text = f"{key.replace('_', ' ').title()}: {nps_data.get(key, 'N/A')}"
        p2.font.size = Pt(18)

    prs.save(output_path)
    return output_path
